"""Extract text from common office file formats.

Supports: PDF, DOCX, XLSX, PPTX, CSV
Saves extracted text as a companion .txt file alongside the original.
"""

import io
from pathlib import Path


# File extensions that need text extraction
EXTRACTABLE_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}

# Text extensions that don't need extraction
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
             ".js", ".ts", ".py", ".html", ".css", ".sh", ".log", ".env"}


def needs_extraction(filename: str) -> bool:
    """Check if a file needs text extraction."""
    ext = Path(filename).suffix.lower()
    return ext in EXTRACTABLE_EXTS


def extract_text(file_bytes: bytes, filename: str) -> str | None:
    """Extract text from a binary file.
    
    Returns extracted text string, or None if extraction fails.
    """
    ext = Path(filename).suffix.lower()
    
    try:
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        elif ext == ".docx":
            return _extract_docx(file_bytes)
        elif ext == ".xlsx":
            return _extract_xlsx(file_bytes)
        elif ext == ".pptx":
            return _extract_pptx(file_bytes)
    except Exception as e:
        print(f"[TextExtractor] Failed to extract from {filename}: {e}")
        return None
    
    return None


def save_extracted_text(save_path: Path, file_bytes: bytes, filename: str) -> Path | None:
    """Extract text and save as a companion .txt file.
    
    For example: report.pdf → report.pdf.txt
    Returns the path to the text file, or None if extraction failed.
    """
    text = extract_text(file_bytes, filename)
    if not text or not text.strip():
        return None
    
    txt_path = save_path.parent / f"{save_path.name}.txt"
    txt_path.write_text(text, encoding="utf-8")
    print(f"[TextExtractor] Extracted {len(text)} chars from {filename} → {txt_path.name}")
    return txt_path


def _extract_pdf(data: bytes) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber
    
    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- 第{i+1}页 ---\n{text.strip()}")
            
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                if table:
                    rows = []
                    for row in table:
                        cells = [str(c or "").strip() for c in row]
                        rows.append(" | ".join(cells))
                    if rows:
                        pages.append("表格:\n" + "\n".join(rows))
    
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    
    doc = Document(io.BytesIO(data))
    parts = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve heading hierarchy
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                parts.append(f"{'#' * level} {text}")
            else:
                parts.append(text)
    
    # Extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n表格:\n" + "\n".join(rows))
    
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        
        if rows:
            parts.append(f"## 工作表: {sheet}\n" + "\n".join(rows))
    
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    
    prs = Presentation(io.BytesIO(data))
    parts = []
    
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        
        if texts:
            parts.append(f"--- 幻灯片 {i+1} ---\n" + "\n".join(texts))
    
    return "\n\n".join(parts)
